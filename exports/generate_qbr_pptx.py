from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# UiPath brand colours
UIPATH_ORANGE = RGBColor(0xFA, 0x46, 0x16)
DARK_BG       = RGBColor(0x1C, 0x1C, 0x2E)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xF5, 0xF5, 0xF7)
MID_GRAY      = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT     = RGBColor(0x1A, 0x1A, 0x2E)
GREEN         = RGBColor(0x2E, 0xCC, 0x71)
YELLOW        = RGBColor(0xF3, 0x9C, 0x12)
RED           = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE_LIGHT  = RGBColor(0xFF, 0xEE, 0xE8)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size=Pt(14), bold=False,
             color=DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def header_bar(slide, title, subtitle=None):
    # Dark top bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill_rgb=DARK_BG)
    # Orange accent line
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Pt(4), fill_rgb=UIPATH_ORANGE)
    add_text(slide, title, Inches(0.4), Inches(0.18), Inches(10), Inches(0.6),
             font_size=Pt(28), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.72), Inches(12), Inches(0.38),
                 font_size=Pt(13), color=MID_GRAY)


def footer(slide, text="Partner Success Program · Q2 FY27 CFO Review · CONFIDENTIAL — UiPath Internal"):
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_rgb=DARK_BG)
    add_text(slide, text, Inches(0.3), Inches(7.2), Inches(12), Inches(0.3),
             font_size=Pt(8), color=MID_GRAY)


def kpi_badge(slide, value, label, x, y, status_color=None):
    w, h = Inches(2.8), Inches(1.3)
    add_rect(slide, x, y, w, h, fill_rgb=LIGHT_GRAY,
             line_rgb=status_color or MID_GRAY, line_width=Pt(2))
    add_text(slide, value, x, y + Inches(0.1), w, Inches(0.7),
             font_size=Pt(30), bold=True, color=status_color or DARK_TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + Inches(0.75), w, Inches(0.45),
             font_size=Pt(11), color=DARK_TEXT, align=PP_ALIGN.CENTER)


# ── SLIDE 1 — COVER ─────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=DARK_BG)
add_rect(slide, 0, Inches(2.8), SLIDE_W, Pt(5), fill_rgb=UIPATH_ORANGE)
add_rect(slide, 0, Inches(4.9), SLIDE_W, Pt(5), fill_rgb=UIPATH_ORANGE)
add_text(slide, "PARTNER SUCCESS PROGRAM",
         Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         font_size=Pt(36), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Q2 FY27 — CFO BUSINESS REVIEW",
         Inches(1), Inches(2.1), Inches(11), Inches(0.65),
         font_size=Pt(24), color=UIPATH_ORANGE, align=PP_ALIGN.CENTER)
add_text(slide, "Presented by: Mihai Gavan, Partner Success Team",
         Inches(1), Inches(3.2), Inches(11), Inches(0.45),
         font_size=Pt(14), color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "June 2026",
         Inches(1), Inches(3.65), Inches(11), Inches(0.4),
         font_size=Pt(14), color=MID_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "CONFIDENTIAL — UiPath Internal",
         Inches(1), Inches(5.2), Inches(11), Inches(0.4),
         font_size=Pt(12), color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 2 — EXECUTIVE SUMMARY ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
header_bar(slide, "Executive Summary", "Q2 FY27 · Partner Success Program")
footer(slide)

# Headline box
add_rect(slide, Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.75),
         fill_rgb=ORANGE_LIGHT, line_rgb=UIPATH_ORANGE, line_width=Pt(2))
add_text(slide, "Every $1 invested in Partner Success returned $4.2 in partner-influenced ARR this quarter",
         Inches(0.55), Inches(1.38), Inches(12.2), Inches(0.65),
         font_size=Pt(14), bold=True, color=UIPATH_ORANGE)

bullets = [
    "Program live since Q1 FY27 — 18 partners enrolled, $2.6M subscription ARR",
    "Enrolled partners outperform non-enrolled: 2.3× more pipeline, 94% vs 81% retention",
    "Gross margin: 68% — ahead of 65% target",
    "Deal velocity: 22% faster for enrolled partners",
]
for i, b in enumerate(bullets):
    add_text(slide, f"•  {b}", Inches(0.5), Inches(2.25) + Inches(0.38) * i,
             Inches(12.3), Inches(0.38), font_size=Pt(13), color=DARK_TEXT)

add_text(slide, "THREE ASKS TODAY", Inches(0.5), Inches(4.0), Inches(4), Inches(0.35),
         font_size=Pt(13), bold=True, color=DARK_BG)
asks = [
    "1.  Approve 8 additional partner slots (H2 FY27)  →  +$800k–$1.2M ARR",
    "2.  Fund 2 additional PSM headcount  →  prevents capacity bottleneck",
    "3.  Unblock Salesforce pipeline attribution (RevOps, 2 wks)  →  accurate Q3 board numbers",
]
for i, a in enumerate(asks):
    add_text(slide, a, Inches(0.5), Inches(4.42) + Inches(0.38) * i,
             Inches(12.3), Inches(0.38), font_size=Pt(13), color=DARK_TEXT)


# ── SLIDE 3 — PROGRAM AT A GLANCE ───────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
header_bar(slide, "Program at a Glance", "Key metrics · Q2 FY27")
footer(slide)

kpi_badge(slide, "$2.6M",  "Subscription ARR",       Inches(0.4),  Inches(1.5),  GREEN)
kpi_badge(slide, "18",     "Partners Enrolled",       Inches(3.4),  Inches(1.5),  GREEN)
kpi_badge(slide, "68%",    "Gross Margin",            Inches(6.4),  Inches(1.5),  GREEN)
kpi_badge(slide, "$10.9M", "Partner-Influenced ARR",  Inches(9.4),  Inches(1.5),  GREEN)

add_text(slide, "12 Lightweight  ·  6 Full Model", Inches(3.4), Inches(2.85),
         Inches(2.8), Inches(0.3), font_size=Pt(10), color=DARK_TEXT, align=PP_ALIGN.CENTER)

# ROI section
add_rect(slide, Inches(0.4), Inches(3.3), Inches(12.5), Inches(1.3), fill_rgb=LIGHT_GRAY)
add_text(slide, "PROGRAM ROI", Inches(0.6), Inches(3.35), Inches(3), Inches(0.35),
         font_size=Pt(13), bold=True, color=DARK_BG)
roi_items = [
    ("Investment (annualised)", "$1.1M"),
    ("Partner-influenced ARR (Q2)", "$10.9M  →  4.2× ROI"),
    ("Partner-sourced new ARR (Q2)", "$3.2M  →  1.9× ROI"),
    ("Projected 12-month ROI", "6.8×"),
]
for i, (label, val) in enumerate(roi_items):
    col = Inches(0.6) if i < 2 else Inches(6.7)
    row = Inches(3.75) if i % 2 == 0 else Inches(4.1)
    if i >= 2:
        row = Inches(3.75) if i == 2 else Inches(4.1)
    add_text(slide, f"{label}:", col, row, Inches(3.5), Inches(0.35),
             font_size=Pt(12), color=DARK_TEXT)
    add_text(slide, val, col + Inches(3.6), row, Inches(2.8), Inches(0.35),
             font_size=Pt(12), bold=True, color=UIPATH_ORANGE)


# ── SLIDE 4 — KPI SCORECARD ─────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
header_bar(slide, "KPI Scorecard", "Q2 FY27 vs FY27 Targets")
footer(slide)

rows = [
    ("Partner Churn Rate",           "<10%",      "5.6%",       "🟢 Ahead",    GREEN),
    ("Partner-Sourced Pipeline",      "+25% YoY",  "+31%",       "🟢 Ahead",    GREEN),
    ("Program Gross Margin",          ">65%",      "68%",        "🟢 On Track", GREEN),
    ("Net Revenue Retention",         "+10% YoY",  "+7.8%",      "🟡 At Risk",  YELLOW),
    ("Technical Health Score",        ">85%",      "82%",        "🟡 At Risk",  YELLOW),
    ("Migration Success Rate",        ">90%",      "87%",        "🟡 At Risk",  YELLOW),
    ("RSI Competency Growth",         "+20% YoY",  "+17%",       "🟡 At Risk",  YELLOW),
    ("QBR Partner NPS",               ">80",       "76",         "🟡 At Risk",  YELLOW),
    ("Agentic Certification Growth",  "+50% YoY",  "+34%",       "🔴 Behind",   RED),
    ("Playbook Adoption (PSMs)",      "95%",       "72%",        "🔴 Behind",   RED),
    ("Partner Portal Logins",         "15k/mo",    "9,200/mo",   "🔴 Behind",   RED),
    ("Artifact Reuse Rate",           "≥60%",      "44%",        "🔴 Behind",   RED),
]

headers = ["KPI", "FY27 Target", "Q2 Actual", "Status"]
col_x   = [Inches(0.3), Inches(5.8), Inches(8.0), Inches(10.1)]
col_w   = [Inches(5.4), Inches(2.1), Inches(2.0), Inches(3.0)]
row_h   = Inches(0.44)
top     = Inches(1.35)

# Header row
add_rect(slide, Inches(0.3), top, Inches(12.7), row_h, fill_rgb=DARK_BG)
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_text(slide, hdr, cx + Inches(0.05), top + Pt(4), cw, row_h,
             font_size=Pt(11), bold=True, color=WHITE)

for i, (kpi, target, actual, status, sc) in enumerate(rows):
    y = top + row_h * (i + 1)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.3), y, Inches(12.7), row_h, fill_rgb=bg)
    add_text(slide, kpi,    col_x[0] + Inches(0.05), y + Pt(3), col_w[0], row_h, font_size=Pt(11), color=DARK_TEXT)
    add_text(slide, target, col_x[1] + Inches(0.05), y + Pt(3), col_w[1], row_h, font_size=Pt(11), color=DARK_TEXT)
    add_text(slide, actual, col_x[2] + Inches(0.05), y + Pt(3), col_w[2], font_size=Pt(11), bold=True, color=sc, h=row_h)
    add_text(slide, status, col_x[3] + Inches(0.05), y + Pt(3), col_w[3], row_h, font_size=Pt(11), color=sc)


# ── SLIDE 5 — FINANCIAL PERFORMANCE ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
header_bar(slide, "Financial Performance", "Subscription Revenue & COGS · Q2 FY27")
footer(slide)

# Revenue table
add_text(slide, "SUBSCRIPTION REVENUE", Inches(0.4), Inches(1.35), Inches(6), Inches(0.35),
         font_size=Pt(12), bold=True, color=DARK_BG)
rev_rows = [
    ("Lightweight (12 partners × $100k)", "$1.2M", "70%"),
    ("Full Model (6 partners × $300k)",   "$1.8M", "70%"),
    ("Y1 Discounts applied",              "-$400k", "—"),
    ("TOTAL",                             "$2.6M",  "68%"),
]
rev_headers = ["Tier", "ARR", "Margin"]
rcol_x = [Inches(0.4), Inches(5.0), Inches(6.2)]
rcol_w = [Inches(4.5), Inches(1.1), Inches(1.0)]
rh = Inches(0.42)

add_rect(slide, Inches(0.4), Inches(1.72), Inches(7.0), rh, fill_rgb=DARK_BG)
for j, (h, cx, cw) in enumerate(zip(rev_headers, rcol_x, rcol_w)):
    add_text(slide, h, cx + Inches(0.05), Inches(1.72) + Pt(4), cw, rh,
             font_size=Pt(11), bold=True, color=WHITE)

for i, (tier, arr, margin) in enumerate(rev_rows):
    y   = Inches(1.72) + rh * (i + 1)
    bg  = LIGHT_GRAY if i % 2 == 0 else WHITE
    bld = i == 3
    if i == 3:
        add_rect(slide, Inches(0.4), y, Inches(7.0), rh, fill_rgb=ORANGE_LIGHT,
                 line_rgb=UIPATH_ORANGE, line_width=Pt(1))
    else:
        add_rect(slide, Inches(0.4), y, Inches(7.0), rh, fill_rgb=bg)
    add_text(slide, tier,   rcol_x[0] + Inches(0.05), y + Pt(3), rcol_w[0], rh, font_size=Pt(11), bold=bld, color=DARK_TEXT)
    add_text(slide, arr,    rcol_x[1] + Inches(0.05), y + Pt(3), rcol_w[1], rh, font_size=Pt(11), bold=bld, color=UIPATH_ORANGE if bld else DARK_TEXT)
    add_text(slide, margin, rcol_x[2] + Inches(0.05), y + Pt(3), rcol_w[2], rh, font_size=Pt(11), bold=bld, color=DARK_TEXT)

# COGS table
add_text(slide, "Q2 COGS BREAKDOWN", Inches(7.8), Inches(1.35), Inches(5), Inches(0.35),
         font_size=Pt(12), bold=True, color=DARK_BG)
cogs_rows = [
    ("PSM headcount (3 FTE)",      "$187k", "41%"),
    ("Solution Architect (0.5)",   "$65k",  "14%"),
    ("Tooling",                    "$38k",   "8%"),
    ("Training & workshops",       "$112k", "25%"),
    ("Travel & engagement",        "$55k",  "12%"),
    ("TOTAL Q2",                   "$457k", "100%"),
]
ccol_x = [Inches(7.8), Inches(11.0), Inches(12.1)]
ccol_w = [Inches(3.1), Inches(1.0),  Inches(1.0)]

add_rect(slide, Inches(7.8), Inches(1.72), Inches(5.2), rh, fill_rgb=DARK_BG)
for j, (h, cx, cw) in enumerate(zip(["Category", "Cost", "%"], ccol_x, ccol_w)):
    add_text(slide, h, cx + Inches(0.05), Inches(1.72) + Pt(4), cw, rh,
             font_size=Pt(11), bold=True, color=WHITE)

for i, (cat, cost, pct) in enumerate(cogs_rows):
    y   = Inches(1.72) + rh * (i + 1)
    bld = i == 5
    bg  = ORANGE_LIGHT if bld else (LIGHT_GRAY if i % 2 == 0 else WHITE)
    add_rect(slide, Inches(7.8), y, Inches(5.2), rh,
             fill_rgb=bg, line_rgb=UIPATH_ORANGE if bld else None, line_width=Pt(1))
    add_text(slide, cat,  ccol_x[0] + Inches(0.05), y + Pt(3), ccol_w[0], rh, font_size=Pt(11), bold=bld, color=DARK_TEXT)
    add_text(slide, cost, ccol_x[1] + Inches(0.05), y + Pt(3), ccol_w[1], rh, font_size=Pt(11), bold=bld, color=UIPATH_ORANGE if bld else DARK_TEXT)
    add_text(slide, pct,  ccol_x[2] + Inches(0.05), y + Pt(3), ccol_w[2], rh, font_size=Pt(11), bold=bld, color=DARK_TEXT)


# ── SLIDE 6 — ENROLLED vs NON-ENROLLED ──────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
header_bar(slide, "Why Partner Success Works", "Enrolled vs Non-Enrolled · Same Tier · Same Region")
footer(slide)

add_text(slide, "Partners in the program consistently outperform peers on every metric that matters to revenue and retention.",
         Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.45), font_size=Pt(13), color=DARK_TEXT)

cmp_rows = [
    ("Annual Retention Rate",       "94%",      "81%",       "+13pp"),
    ("Partner-Sourced Pipeline",    "$1.8M avg", "$0.78M avg", "+2.3×"),
    ("Deal Velocity",               "47 days",  "60 days",   "-22%"),
    ("Agentic Cert Rate",           "68%",      "29%",       "+39pp"),
    ("Escalations per Quarter",     "0.8",      "2.1",       "-62%"),
]
cmp_headers = ["Metric", "Enrolled", "Non-Enrolled", "Delta"]
ccol_x = [Inches(0.4), Inches(5.5), Inches(8.3), Inches(11.0)]
ccol_w = [Inches(5.0), Inches(2.7), Inches(2.6), Inches(2.2)]
rh = Inches(0.56)
top = Inches(1.95)

add_rect(slide, Inches(0.4), top, Inches(12.5), rh, fill_rgb=DARK_BG)
for j, (h, cx, cw) in enumerate(zip(cmp_headers, ccol_x, ccol_w)):
    add_text(slide, h, cx + Inches(0.08), top + Pt(6), cw, rh,
             font_size=Pt(12), bold=True, color=WHITE)

for i, (metric, enrolled, non_enr, delta) in enumerate(cmp_rows):
    y  = top + rh * (i + 1)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.4), y, Inches(12.5), rh, fill_rgb=bg)
    add_text(slide, metric,   ccol_x[0] + Inches(0.08), y + Pt(6), ccol_w[0], rh, font_size=Pt(12), color=DARK_TEXT)
    add_text(slide, enrolled, ccol_x[1] + Inches(0.08), y + Pt(6), ccol_w[1], rh, font_size=Pt(12), bold=True, color=GREEN)
    add_text(slide, non_enr,  ccol_x[2] + Inches(0.08), y + Pt(6), ccol_w[2], rh, font_size=Pt(12), color=RED)
    add_text(slide, delta,    ccol_x[3] + Inches(0.08), y + Pt(6), ccol_w[3], rh, font_size=Pt(12), bold=True, color=UIPATH_ORANGE)


# ── SLIDE 7 — STRATEGIC ALIGNMENT ───────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
header_bar(slide, "Strategic Alignment", "FY28/29 Commercial Motions")
footer(slide)

motions = [
    ("Motion 1", "Install-Base Extension", "CoE → LOB at scale",
     ["8 enrolled partners managing $47M in UiPath ARR",
      "Target: 12 partners running LOB expansion plays by FY27 close",
      "Gap: Only 3 of 8 have signed co-sell motion agreements with AEs"]),
    ("Motion 2", "Vertical Solutions", "Commercial-segment economics",
     ["4 partners with published vertical packages (FinCrime, RCM, P2P, O2C)",
      "Accenture: $4.2M healthcare vertical deal closed Q2",
      "Gap: BOAT/vertical language adoption at 38% — target 80%"]),
    ("Motion 3", "Test Cloud / Dark Factory", "Autonomous testing",
     ["Accenture, Deloitte, Capgemini in active re-skilling programs",
      "2 of 3 committed to Test Cloud pilots Q3 FY27",
      "Risk: Capgemini test lead vacancy — escalated to exec sponsor"]),
]

for i, (tag, title, sub, points) in enumerate(motions):
    x   = Inches(0.4) + Inches(4.3) * i
    bx_y = Inches(1.35)
    add_rect(slide, x, bx_y, Inches(4.1), Inches(5.55),
             fill_rgb=LIGHT_GRAY, line_rgb=UIPATH_ORANGE, line_width=Pt(2))
    add_rect(slide, x, bx_y, Inches(4.1), Inches(0.75), fill_rgb=UIPATH_ORANGE)
    add_text(slide, tag,   x + Inches(0.12), bx_y + Pt(2), Inches(3.9), Inches(0.3),
             font_size=Pt(10), bold=True, color=WHITE)
    add_text(slide, title, x + Inches(0.12), bx_y + Inches(0.28), Inches(3.9), Inches(0.38),
             font_size=Pt(14), bold=True, color=WHITE)
    add_text(slide, sub,   x + Inches(0.12), bx_y + Inches(0.75), Inches(3.9), Inches(0.35),
             font_size=Pt(10), color=DARK_TEXT)
    for j, pt in enumerate(points):
        add_text(slide, f"•  {pt}", x + Inches(0.12),
                 bx_y + Inches(1.2) + Inches(0.75) * j,
                 Inches(3.85), Inches(0.7), font_size=Pt(11), color=DARK_TEXT)


# ── SLIDE 8 — TOP WINS ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
header_bar(slide, "Top Wins — Q2 FY27", "Partner-driven outcomes")
footer(slide)

wins = [
    ("Heineken Global Rollout",          "Roboyo",      "$1.1M new ARR · 3-year deal"),
    ("Accenture FinCrime Vertical",       "Accenture",   "$4.2M influenced pipeline · 6 new logos"),
    ("Deloitte Agentic POC Launched",     "Deloitte",    "Manufacturing · Expansion play Q3"),
    ("Greenlight — 0 P1 Escalations",     "Greenlight",  "2nd consecutive quarter · NPS 94"),
    ("14 New Agentic Certifications",     "Multiple",    "+34% YoY cert growth"),
]
w_headers = ["Win", "Partner", "Impact"]
wcol_x = [Inches(0.4), Inches(7.5), Inches(9.5)]
wcol_w = [Inches(7.0), Inches(1.9), Inches(3.7)]
rh = Inches(0.75)
top = Inches(1.35)

add_rect(slide, Inches(0.4), top, Inches(12.7), Inches(0.45), fill_rgb=DARK_BG)
for j, (h, cx, cw) in enumerate(zip(w_headers, wcol_x, wcol_w)):
    add_text(slide, h, cx + Inches(0.08), top + Pt(4), cw, Inches(0.45),
             font_size=Pt(11), bold=True, color=WHITE)

for i, (win, partner, impact) in enumerate(wins):
    y  = top + Inches(0.45) + rh * i
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.4), y, Inches(12.7), rh, fill_rgb=bg)
    add_rect(slide, Inches(0.4), y, Inches(0.08), rh, fill_rgb=GREEN)
    add_text(slide, win,     wcol_x[0] + Inches(0.15), y + Pt(6), wcol_w[0], rh, font_size=Pt(12), bold=True, color=DARK_TEXT)
    add_text(slide, partner, wcol_x[1] + Inches(0.08), y + Pt(6), wcol_w[1], rh, font_size=Pt(11), color=DARK_TEXT)
    add_text(slide, impact,  wcol_x[2] + Inches(0.08), y + Pt(6), wcol_w[2], rh, font_size=Pt(11), color=UIPATH_ORANGE)


# ── SLIDE 9 — RISKS ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
header_bar(slide, "Risks & Mitigations", "Q2 FY27 Risk Register")
footer(slide)

risk_rows = [
    ("PSM bottleneck at 20+ partners", "High",   "High",   "Needs headcount approval",        RED),
    ("NRR below target (2 partners)",  "Medium", "High",   "Recovery plans in progress",      YELLOW),
    ("Agentic cert growth behind",     "High",   "Medium", "Boot camp Aug 2026 scheduled",    YELLOW),
    ("Salesforce attribution gap",     "High",   "High",   "CFO unblock needed  ← KEY ASK",  RED),
    ("Capgemini test lead vacancy",    "Low",    "Medium", "Escalated to exec sponsor",       GREEN),
]
r_headers = ["Risk", "Probability", "Impact", "Mitigation"]
rcol_x = [Inches(0.4), Inches(6.5), Inches(8.2), Inches(10.0)]
rcol_w = [Inches(6.0), Inches(1.6), Inches(1.7), Inches(3.2)]
rh = Inches(0.72)
top = Inches(1.35)

add_rect(slide, Inches(0.4), top, Inches(12.7), Inches(0.45), fill_rgb=DARK_BG)
for j, (h, cx, cw) in enumerate(zip(r_headers, rcol_x, rcol_w)):
    add_text(slide, h, cx + Inches(0.08), top + Pt(4), cw, Inches(0.45),
             font_size=Pt(11), bold=True, color=WHITE)

for i, (risk, prob, impact, mit, sc) in enumerate(risk_rows):
    y  = top + Inches(0.45) + rh * i
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.4), y, Inches(12.7), rh, fill_rgb=bg)
    add_rect(slide, Inches(0.4), y, Inches(0.1), rh, fill_rgb=sc)
    add_text(slide, risk,   rcol_x[0] + Inches(0.15), y + Pt(6), rcol_w[0], rh, font_size=Pt(11), bold=True, color=DARK_TEXT)
    add_text(slide, prob,   rcol_x[1] + Inches(0.08), y + Pt(6), rcol_w[1], rh, font_size=Pt(11), color=sc)
    add_text(slide, impact, rcol_x[2] + Inches(0.08), y + Pt(6), rcol_w[2], rh, font_size=Pt(11), color=sc)
    add_text(slide, mit,    rcol_x[3] + Inches(0.08), y + Pt(6), rcol_w[3], rh, font_size=Pt(11), color=DARK_TEXT)


# ── SLIDE 10 — H2 FY27 ASKS ──────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
header_bar(slide, "H2 FY27 — Three Asks", "Decisions needed today")
footer(slide)

asks_data = [
    ("ASK 1", "Expand to 26 Partners (+8 Slots)",
     ["11 qualified partners in pipeline, ready to enroll",
      "Incremental ARR: +$800k–$1.2M  |  COGS: +$240k–$360k",
      "Net margin uplift: +$440k–$840k annualised"]),
    ("ASK 2", "2 Additional PSM Headcount",
     ["Cost: $260k fully loaded",
      "Unlocks: capacity for 26–30 partners without quality risk",
      "Risk if not approved: bottleneck at 20, churn risk on existing base"]),
    ("ASK 3", "Salesforce Attribution — RevOps Unblock",
     ["Issue: influenced pipeline undercounted by ~35%",
      "Impact: Q3 board numbers will understate program ROI",
      "Ask: CFO directive to RevOps · Deadline: July 31 · Effort: 2 wks"]),
]

for i, (tag, title, points) in enumerate(asks_data):
    x  = Inches(0.4) + Inches(4.3) * i
    by = Inches(1.35)
    add_rect(slide, x, by, Inches(4.1), Inches(5.5),
             fill_rgb=LIGHT_GRAY, line_rgb=UIPATH_ORANGE, line_width=Pt(2))
    add_rect(slide, x, by, Inches(4.1), Inches(0.75), fill_rgb=DARK_BG)
    add_text(slide, tag,   x + Inches(0.15), by + Pt(4), Inches(3.9), Inches(0.3),
             font_size=Pt(11), bold=True, color=UIPATH_ORANGE)
    add_text(slide, title, x + Inches(0.15), by + Inches(0.32), Inches(3.9), Inches(0.38),
             font_size=Pt(13), bold=True, color=WHITE)
    for j, pt in enumerate(points):
        add_text(slide, f"•  {pt}", x + Inches(0.15),
                 by + Inches(0.95) + Inches(0.85) * j,
                 Inches(3.82), Inches(0.82), font_size=Pt(11), color=DARK_TEXT)


# ── SLIDE 11 — NEXT STEPS ────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
header_bar(slide, "Next Steps", "Actions and owners coming out of today")
footer(slide)

ns_rows = [
    ("Decision: H2 expansion + 2 PSM hires",      "CFO",              "Jun 25, 2026", RED),
    ("RevOps Salesforce attribution kickoff",       "CFO → RevOps",     "Jul 1, 2026",  RED),
    ("Agentic boot camp invitations sent",          "Enablement team",  "Jun 20, 2026", YELLOW),
    ("Recovery plans for 2 underperforming PSMs",  "PSM team",         "Jun 18, 2026", YELLOW),
    ("Next CFO review (Q3 FY27 close)",            "CFO + VP PS",      "Sep 10, 2026", GREEN),
]
ns_headers = ["Action", "Owner", "Due Date"]
ncol_x = [Inches(0.4), Inches(7.5), Inches(10.5)]
ncol_w = [Inches(7.0), Inches(2.9), Inches(2.7)]
rh = Inches(0.75)
top = Inches(1.5)

add_rect(slide, Inches(0.4), top, Inches(12.7), Inches(0.45), fill_rgb=DARK_BG)
for j, (h, cx, cw) in enumerate(zip(ns_headers, ncol_x, ncol_w)):
    add_text(slide, h, cx + Inches(0.08), top + Pt(4), cw, Inches(0.45),
             font_size=Pt(11), bold=True, color=WHITE)

for i, (action, owner, date, sc) in enumerate(ns_rows):
    y  = top + Inches(0.45) + rh * i
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.4), y, Inches(12.7), rh, fill_rgb=bg)
    add_rect(slide, Inches(0.4), y, Inches(0.1), rh, fill_rgb=sc)
    add_text(slide, action, ncol_x[0] + Inches(0.15), y + Pt(10), ncol_w[0], rh, font_size=Pt(12), bold=True, color=DARK_TEXT)
    add_text(slide, owner,  ncol_x[1] + Inches(0.08), y + Pt(10), ncol_w[1], rh, font_size=Pt(12), color=DARK_TEXT)
    add_text(slide, date,   ncol_x[2] + Inches(0.08), y + Pt(10), ncol_w[2], rh, font_size=Pt(12), bold=True, color=sc)


# ── SAVE ─────────────────────────────────────────────────────────────────────
out = "/Users/mihai.gavan/Documents/ClaudeCode-PB/exports/QBR_CFO_Review_Q2_FY27.pptx"
prs.save(out)
print(f"Saved: {out}")
