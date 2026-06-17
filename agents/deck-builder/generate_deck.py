"""
UiPath Partner Success — Deck Builder Engine
Reads agents/deck-builder/current_spec.json → writes exports/<filename>.pptx
Usage: python3 agents/deck-builder/generate_deck.py
"""

import json, sys, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── UIPATH 2025 OFFICIAL BRAND PALETTE ───────────────────────────────────────
# Source: 260220-UiPath-PPT-2026-RELEASE.pptx theme "UiPath 2025"
# Font: Arial (official presentation font per PPTX master template)

T = {
    # Primary brand colors
    "orange":      RGBColor(0xFA, 0x46, 0x16),  # Robotic Orange — accent1, primary CTA
    "dark":        RGBColor(0x18, 0x21, 0x26),  # Deep Blue — dk2, slide backgrounds
    "teal":        RGBColor(0x0B, 0xA2, 0xB3),  # Agentic Teal — accent3
    # Secondary brand colors
    "dark_blue":   RGBColor(0x1E, 0x64, 0x82),  # Dark Blue — accent2
    "mid_blue":    RGBColor(0x5B, 0xCB, 0xDE),  # Mid Blue
    "yellow":      RGBColor(0xDA, 0x91, 0x00),  # Yellow — accent4
    "magenta":     RGBColor(0xCC, 0x19, 0x56),  # Magenta — accent5
    "purple":      RGBColor(0x8B, 0x28, 0x8A),  # Testing Purple — accent6
    "green":       RGBColor(0x99, 0xC2, 0x57),  # Green
    # Neutrals
    "white":       RGBColor(0xFF, 0xFF, 0xFF),  # White — lt1
    "black":       RGBColor(0x00, 0x00, 0x00),  # Black — dk1
    "dark_grey":   RGBColor(0x61, 0x61, 0x61),  # Dark Grey
    "darker_grey": RGBColor(0x48, 0x48, 0x48),  # Darker Grey
    "offset_grey": RGBColor(0xD9, 0xD9, 0xD9),  # Offset Grey — lt2
    "light_grey":  RGBColor(0xF6, 0xF6, 0xF6),  # Light Grey
    # Block (darker shade) variants
    "block_orange":    RGBColor(0xCE, 0x34, 0x0B),  # Block Orange
    "block_teal":      RGBColor(0x15, 0x83, 0x9A),  # Block Teal
    "block_deep_blue": RGBColor(0x2D, 0x37, 0x3C),  # Block Deep Blue
    "slate_blue":      RGBColor(0x43, 0x5D, 0x6B),  # Slate Blue
    # Aliases used by slide builders
    "dark_card":   RGBColor(0x2D, 0x37, 0x3C),  # Block Deep Blue (card surfaces)
    "off_white":   RGBColor(0xF6, 0xF6, 0xF6),  # Light Grey (off-white surfaces)
    "mid_gray":    RGBColor(0x61, 0x61, 0x61),  # Dark Grey (secondary text)
    "light_gray":  RGBColor(0xD9, 0xD9, 0xD9),  # Offset Grey (borders)
    "text":        RGBColor(0x18, 0x21, 0x26),  # Deep Blue (body text on light bg)
    "red":         RGBColor(0xCC, 0x19, 0x56),  # Magenta used as red/danger
    "blue":        RGBColor(0x1E, 0x64, 0x82),  # Dark Blue
    "orange_bg":   RGBColor(0xFF, 0xEE, 0xE8),  # Orange tint background
}

COLOR_MAP = {
    "green":   T["green"],
    "orange":  T["orange"],
    "blue":    T["blue"],
    "red":     T["red"],
    "default": T["orange"],
}

W = Inches(13.33)
H = Inches(7.5)

# ── HELPERS ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def box(slide, x, y, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=Pt(13), bold=False,
        color=None, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    color = color or T["text"]
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Arial"  # UiPath official presentation font
    return tb

def header(slide, title, subtitle=None):
    box(slide, 0, 0, W, Inches(1.1), fill=T["dark"])
    box(slide, 0, Inches(1.1), W, Pt(4), fill=T["orange"])
    txt(slide, title, Inches(0.45), Inches(0.15), Inches(10), Inches(0.6),
        size=Pt(26), bold=True, color=T["white"])
    if subtitle:
        txt(slide, subtitle, Inches(0.45), Inches(0.7), Inches(11), Inches(0.35),
            size=Pt(12), color=T["mid_gray"])

def footer(slide, meta):
    box(slide, 0, Inches(7.18), W, Inches(0.32), fill=T["dark"])
    label = f"{meta.get('deck_type','')}  ·  {meta.get('title','')}  ·  {meta.get('date','')}  ·  CONFIDENTIAL — UiPath Internal"
    txt(slide, label, Inches(0.35), Inches(7.19), Inches(12.5), Inches(0.3),
        size=Pt(8), color=T["mid_gray"])

# ── SLIDE BUILDERS ────────────────────────────────────────────────────────────

def build_cover(prs, spec, meta):
    slide = blank_slide(prs)
    box(slide, 0, 0, W, H, fill=T["dark"])
    box(slide, 0, Inches(3.0), W, Pt(5), fill=T["orange"])
    box(slide, 0, Inches(5.1), W, Pt(5), fill=T["orange"])
    # Subtle grid texture
    for i in range(0, 14):
        box(slide, Inches(i), 0, Pt(0.5), H, fill=RGBColor(0x28, 0x28, 0x40))

    txt(slide, meta.get("deck_type", "").upper(),
        Inches(1), Inches(0.9), Inches(11), Inches(0.4),
        size=Pt(13), bold=True, color=T["orange"], align=PP_ALIGN.CENTER)
    txt(slide, meta["title"],
        Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.1),
        size=Pt(44), bold=True, color=T["white"], align=PP_ALIGN.CENTER)
    if meta.get("subtitle"):
        txt(slide, meta["subtitle"],
            Inches(0.5), Inches(2.6), Inches(12.3), Inches(0.6),
            size=Pt(22), color=T["orange"], align=PP_ALIGN.CENTER)
    box(slide, Inches(4.5), Inches(3.2), Inches(4.3), Pt(2), fill=T["mid_gray"])
    if meta.get("presenter"):
        txt(slide, meta["presenter"],
            Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.45),
            size=Pt(15), color=T["white"], align=PP_ALIGN.CENTER)
    if meta.get("date"):
        txt(slide, meta["date"],
            Inches(0.5), Inches(3.95), Inches(12.3), Inches(0.4),
            size=Pt(13), color=T["mid_gray"], align=PP_ALIGN.CENTER)


def build_agenda(prs, spec, meta):
    slide = blank_slide(prs)
    box(slide, 0, 0, W, H, fill=T["white"])
    header(slide, spec.get("title", "Agenda"))
    footer(slide, meta)
    items = spec.get("items", [])
    cols = 2 if len(items) > 5 else 1
    chunk = (len(items) + 1) // 2 if cols == 2 else len(items)
    for idx, item in enumerate(items):
        col = idx // chunk
        row = idx %  chunk
        x = Inches(0.5) + Inches(6.4) * col
        y = Inches(1.5) + Inches(0.82) * row
        box(slide, x, y, Inches(0.55), Inches(0.55), fill=T["orange"])
        txt(slide, str(idx + 1), x, y, Inches(0.55), Inches(0.55),
            size=Pt(18), bold=True, color=T["white"], align=PP_ALIGN.CENTER)
        txt(slide, item, x + Inches(0.65), y + Inches(0.08),
            Inches(5.6), Inches(0.55), size=Pt(15), bold=True, color=T["text"])


def build_section(prs, spec, meta):
    slide = blank_slide(prs)
    box(slide, 0, 0, W, H, fill=T["dark"])
    box(slide, 0, 0, Inches(0.18), H, fill=T["orange"])
    for i in range(0, 14):
        box(slide, Inches(i), 0, Pt(0.5), H, fill=RGBColor(0x28, 0x28, 0x40))
    num = spec.get("number", "")
    if num:
        txt(slide, num, Inches(0.4), Inches(1.8), Inches(3), Inches(1.8),
            size=Pt(96), bold=True, color=RGBColor(0x30, 0x30, 0x50), align=PP_ALIGN.LEFT)
    txt(slide, spec.get("title", ""),
        Inches(0.5), Inches(2.8), Inches(12), Inches(0.9),
        size=Pt(38), bold=True, color=T["white"])
    if spec.get("subtitle"):
        box(slide, Inches(0.5), Inches(3.8), Inches(1.8), Pt(3), fill=T["orange"])
        txt(slide, spec["subtitle"], Inches(0.5), Inches(3.95), Inches(10), Inches(0.4),
            size=Pt(16), color=T["mid_gray"])
    footer(slide, meta)


def build_bullets(prs, spec, meta):
    slide = blank_slide(prs)
    box(slide, 0, 0, W, H, fill=T["white"])
    header(slide, spec.get("title", ""), spec.get("subtitle"))
    footer(slide, meta)
    bullets = spec.get("bullets", [])
    top = Inches(1.35)
    for i, b in enumerate(bullets[:10]):
        y = top + Inches(0.58) * i
        box(slide, Inches(0.45), y + Inches(0.14), Inches(0.22), Inches(0.22), fill=T["orange"])
        txt(slide, b, Inches(0.82), y, Inches(11.9), Inches(0.55),
            size=Pt(14), color=T["text"])


def build_two_col(prs, spec, meta):
    slide = blank_slide(prs)
    box(slide, 0, 0, W, H, fill=T["white"])
    header(slide, spec.get("title", ""))
    footer(slide, meta)
    left  = spec.get("left",  {})
    right = spec.get("right", {})
    for col_idx, col in enumerate([left, right]):
        x = Inches(0.4) + Inches(6.5) * col_idx
        box(slide, x, Inches(1.3), Inches(6.2), Inches(5.65),
            fill=T["off_white"], line=T["light_gray"], lw=Pt(1))
        box(slide, x, Inches(1.3), Inches(6.2), Inches(0.5), fill=T["dark"])
        txt(slide, col.get("heading", ""), x + Inches(0.15), Inches(1.33),
            Inches(5.9), Inches(0.45), size=Pt(14), bold=True, color=T["white"])
        for i, b in enumerate(col.get("bullets", [])[:8]):
            y = Inches(1.95) + Inches(0.58) * i
            box(slide, x + Inches(0.15), y + Inches(0.14),
                Inches(0.18), Inches(0.18), fill=T["orange"])
            txt(slide, b, x + Inches(0.44), y, Inches(5.6), Inches(0.55),
                size=Pt(13), color=T["text"])


def build_stat_row(prs, spec, meta):
    slide = blank_slide(prs)
    box(slide, 0, 0, W, H, fill=T["white"])
    header(slide, spec.get("title", ""))
    footer(slide, meta)
    stats = spec.get("stats", [])[:4]
    n    = len(stats)
    bw   = (Inches(12.5)) / n
    pad  = Inches(0.2)
    for i, s in enumerate(stats):
        x = Inches(0.4) + bw * i
        c = COLOR_MAP.get(s.get("color", "default"), T["orange"])
        box(slide, x + pad, Inches(1.5), bw - pad * 2, Inches(4.2),
            fill=T["off_white"], line=c, lw=Pt(2))
        box(slide, x + pad, Inches(1.5), bw - pad * 2, Pt(4), fill=c)
        txt(slide, s.get("value", ""),
            x + pad, Inches(2.1), bw - pad * 2, Inches(1.4),
            size=Pt(52), bold=True, color=c, align=PP_ALIGN.CENTER)
        txt(slide, s.get("label", ""),
            x + pad, Inches(3.55), bw - pad * 2, Inches(0.5),
            size=Pt(14), color=T["text"], align=PP_ALIGN.CENTER)
        if s.get("sub"):
            txt(slide, s["sub"], x + pad, Inches(4.1), bw - pad * 2, Inches(0.4),
                size=Pt(11), color=T["mid_gray"], align=PP_ALIGN.CENTER)


def build_table(prs, spec, meta):
    slide = blank_slide(prs)
    box(slide, 0, 0, W, H, fill=T["white"])
    header(slide, spec.get("title", ""))
    footer(slide, meta)
    headers = spec.get("headers", [])
    rows    = spec.get("rows", [])
    if not headers: return
    n_cols  = len(headers)
    col_w   = Inches(12.5) / n_cols
    row_h   = Inches(0.48)
    top     = Inches(1.3)
    box(slide, Inches(0.4), top, Inches(12.5), row_h, fill=T["dark"])
    for j, h in enumerate(headers):
        txt(slide, h, Inches(0.45) + col_w * j, top + Pt(4),
            col_w - Inches(0.1), row_h,
            size=Pt(11), bold=True, color=T["white"])
    for i, row in enumerate(rows[:14]):
        y  = top + row_h * (i + 1)
        bg = T["off_white"] if i % 2 == 0 else T["white"]
        box(slide, Inches(0.4), y, Inches(12.5), row_h, fill=bg)
        for j, cell in enumerate(row[:n_cols]):
            cell_str = str(cell)
            cell_color = T["text"]
            if any(x in cell_str for x in ["🟢", "✅", "Ahead", "Done"]):
                cell_color = T["green"]
            elif any(x in cell_str for x in ["🔴", "❌", "Behind", "Off"]):
                cell_color = T["red"]
            elif any(x in cell_str for x in ["🟡", "⚠️", "At Risk"]):
                cell_color = T["yellow"]
            txt(slide, cell_str, Inches(0.45) + col_w * j, y + Pt(3),
                col_w - Inches(0.1), row_h,
                size=Pt(11), color=cell_color,
                bold=(j == 0))


def build_cards(prs, spec, meta):
    slide = blank_slide(prs)
    box(slide, 0, 0, W, H, fill=T["white"])
    header(slide, spec.get("title", ""))
    footer(slide, meta)
    cards = spec.get("cards", [])[:4]
    n    = len(cards)
    cw   = Inches(12.5) / n
    pad  = Inches(0.2)
    for i, c in enumerate(cards):
        x = Inches(0.4) + cw * i
        box(slide, x + pad, Inches(1.35), cw - pad * 2, Inches(5.5),
            fill=T["off_white"], line=T["light_gray"], lw=Pt(1))
        box(slide, x + pad, Inches(1.35), cw - pad * 2, Pt(4), fill=T["orange"])
        txt(slide, c.get("icon", ""), x + pad, Inches(1.5),
            cw - pad * 2, Inches(0.8),
            size=Pt(32), align=PP_ALIGN.CENTER)
        txt(slide, c.get("title", ""), x + pad, Inches(2.4),
            cw - pad * 2, Inches(0.5),
            size=Pt(14), bold=True, color=T["text"], align=PP_ALIGN.CENTER)
        box(slide, x + pad + Inches(0.6), Inches(2.95),
            cw - pad * 2 - Inches(1.2), Pt(2), fill=T["orange"])
        txt(slide, c.get("body", ""), x + pad + Inches(0.1), Inches(3.1),
            cw - pad * 2 - Inches(0.2), Inches(3.4),
            size=Pt(12), color=T["text"], align=PP_ALIGN.CENTER)


def build_timeline(prs, spec, meta):
    slide = blank_slide(prs)
    box(slide, 0, 0, W, H, fill=T["white"])
    header(slide, spec.get("title", ""))
    footer(slide, meta)
    steps  = spec.get("steps", [])[:5]
    n      = len(steps)
    if not n: return
    sw     = Inches(12.5) / n
    pad    = Inches(0.25)
    mid_y  = Inches(3.8)
    line_y = mid_y + Inches(0.18)
    box(slide, Inches(0.4), line_y, Inches(12.5), Pt(3), fill=T["mid_gray"])
    status_colors = {"done": T["green"], "active": T["orange"], "upcoming": T["mid_gray"]}
    for i, step in enumerate(steps):
        x  = Inches(0.4) + sw * i
        cx = x + sw / 2
        sc = status_colors.get(step.get("status", "upcoming"), T["mid_gray"])
        box(slide, cx - Inches(0.28), mid_y, Inches(0.55), Inches(0.55),
            fill=sc, line=T["white"], lw=Pt(2))
        if step.get("status") == "done":
            txt(slide, "✓", cx - Inches(0.28), mid_y, Inches(0.55), Inches(0.55),
                size=Pt(14), bold=True, color=T["white"], align=PP_ALIGN.CENTER)
        txt(slide, step.get("phase", ""), cx - sw/2 + pad, Inches(1.5),
            sw - pad * 2, Inches(0.35), size=Pt(11), bold=True,
            color=sc, align=PP_ALIGN.CENTER)
        txt(slide, step.get("title", ""), cx - sw/2 + pad, Inches(1.9),
            sw - pad * 2, Inches(0.5), size=Pt(14), bold=True,
            color=T["text"], align=PP_ALIGN.CENTER)
        txt(slide, step.get("dates", ""), cx - sw/2 + pad, Inches(4.55),
            sw - pad * 2, Inches(0.35), size=Pt(11),
            color=T["mid_gray"], align=PP_ALIGN.CENTER)
        if step.get("detail"):
            txt(slide, step["detail"], cx - sw/2 + pad, Inches(2.5),
                sw - pad * 2, Inches(1.1), size=Pt(11),
                color=T["text"], align=PP_ALIGN.CENTER)


def build_closing(prs, spec, meta):
    slide = blank_slide(prs)
    box(slide, 0, 0, W, H, fill=T["dark"])
    box(slide, 0, 0, Inches(0.18), H, fill=T["orange"])
    for i in range(0, 14):
        box(slide, Inches(i), 0, Pt(0.5), H, fill=RGBColor(0x28, 0x28, 0x40))
    txt(slide, spec.get("title", "Thank You"),
        Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.2),
        size=Pt(52), bold=True, color=T["white"], align=PP_ALIGN.CENTER)
    if spec.get("message"):
        box(slide, Inches(4.5), Inches(3.2), Inches(4.3), Pt(3), fill=T["orange"])
        txt(slide, spec["message"], Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.6),
            size=Pt(20), color=T["orange"], align=PP_ALIGN.CENTER)
    if spec.get("contact"):
        txt(slide, spec["contact"], Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.4),
            size=Pt(14), color=T["mid_gray"], align=PP_ALIGN.CENTER)
    txt(slide, meta.get("date", ""), Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.35),
        size=Pt(11), color=T["mid_gray"], align=PP_ALIGN.CENTER)


# ── BUILDER REGISTRY ─────────────────────────────────────────────────────────

BUILDERS = {
    "cover":    build_cover,
    "agenda":   build_agenda,
    "section":  build_section,
    "bullets":  build_bullets,
    "two_col":  build_two_col,
    "stat_row": build_stat_row,
    "table":    build_table,
    "cards":    build_cards,
    "timeline": build_timeline,
    "closing":  build_closing,
}

# ── MAIN ─────────────────────────────────────────────────────────────────────

def main():
    spec_path = os.path.join(os.path.dirname(__file__), "current_spec.json")
    with open(spec_path) as f:
        spec = json.load(f)

    meta = {k: spec.get(k, "") for k in ["title", "subtitle", "presenter", "date", "deck_type"]}
    prs  = new_prs()

    built = 0
    skipped = []
    for slide_spec in spec.get("slides", []):
        stype = slide_spec.get("type", "")
        if stype in BUILDERS:
            BUILDERS[stype](prs, slide_spec, meta)
            built += 1
        else:
            skipped.append(stype)

    filename = spec.get("output_filename", "deck_output.pptx")
    out_dir  = os.path.join(os.path.dirname(__file__), "../../exports")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, filename)
    prs.save(out_path)

    abs_path = os.path.abspath(out_path)
    print(f"✅  {built} slides built → {abs_path}")
    if skipped:
        print(f"⚠️  Unknown slide types skipped: {skipped}")

if __name__ == "__main__":
    main()
